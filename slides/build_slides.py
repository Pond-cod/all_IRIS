# -*- coding: utf-8 -*-
"""Build Day-3 teaching slides (IRIS BrighterBee branded background)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

_HERE = os.path.dirname(os.path.abspath(__file__))
BG = os.path.join(_HERE, "assets", "bg_iris.jpg")   # IRIS BrighterBee background
OUT = os.path.join(_HERE, "Day3_Workshop_Slides.pptx")
COVER = os.path.join(_HERE, "assets", "bg_cover_dark.jpg")   # dark cover (slide 1)
REPO_URL = "https://github.com/amornpan/advanced-de-applied-analytics"   # clickable repo link
PULL_URL = ("https://jupyter-hub.minddatatech.com/hub/user-redirect/git-pull"
            "?repo=https%3A%2F%2Fgithub.com%2Famornpan%2Fadvanced-de-applied-analytics&branch=main")   # nbgitpuller one-click (user-redirect spawns the server first)

# palette
NAVY  = RGBColor(0x14, 0x30, 0x5F)
BLUE  = RGBColor(0x1F, 0x5F, 0xC7)
CYAN  = RGBColor(0x14, 0x8C, 0xC9)
INK   = RGBColor(0x2A, 0x38, 0x49)
MUTED = RGBColor(0x6B, 0x7A, 0x8D)
GREEN = RGBColor(0x1E, 0x9E, 0x6A)
ORANGE= RGBColor(0xDE, 0x7A, 0x12)
PURPLE= RGBColor(0x6B, 0x3F, 0xA0)
CARD  = RGBColor(0xEE, 0xF3, 0xFB)
CARDLN= RGBColor(0xCF, 0xDE, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG= RGBColor(0x0F, 0x20, 0x3A)
CODETX= RGBColor(0xE8, 0xEE, 0xF7)
CYANL = RGBColor(0x5B, 0xC8, 0xF5)
LIGHT = RGBColor(0xCF, 0xDD, 0xF2)
DQ_BLUE = RGBColor(0x9D, 0xC3, 0xE6)   # dq_power-bi section colours
DQ_YEL  = RGBColor(0xFF, 0xE6, 0x99)
DQ_GRN  = RGBColor(0xC6, 0xE0, 0xB4)

HEAD = "Inter"                 # macOS-like (SF Pro substitute) for Latin
BODY = "Inter"
CODE = "Cascadia Mono"
THAI = "Kanit"                  # complex-script (Thai), Sukhumvit-like — set per run

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
    return s


def _set_para(p, runs, align=PP_ALIGN.LEFT, sa=4, sb=0, lh=None):
    p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if sb is not None: p.space_before = Pt(sb)
    if lh is not None: p.line_spacing = lh
    for (t, sz, b, c, f) in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.name = f
        r.font.color.rgb = c
        # set Thai (complex-script) typeface so Thai glyphs use IBM Plex Sans Thai
        rPr = r._r.get_or_add_rPr()
        cs = rPr.find(qn("a:cs"))
        if cs is None:
            cs = rPr.makeelement(qn("a:cs"), {"typeface": THAI}); rPr.append(cs)
        else:
            cs.set("typeface", THAI)
    return p


def text(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dict(runs=[(t,sz,b,c,f)], align, sa, sb, lh)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, pa["runs"], pa.get("align", PP_ALIGN.LEFT),
                  pa.get("sa", 4), pa.get("sb", 0), pa.get("lh"))
    return tb


def rrect(s, x, y, w, h, fill, line=None, lw=1.0, radius=0.09, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.text_frame.paragraphs[0].text = ""
    return shp


def rect(s, x, y, w, h, fill, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def circle(s, x, y, d, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def badge(s, x, y, d, n, fill=BLUE, tc=WHITE, sz=15):
    circle(s, x, y, d, fill)
    text(s, x, y - 0.02, d, d, [{"runs": [(str(n), sz, True, tc, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}],
         anchor=MSO_ANCHOR.MIDDLE)


def title(s, t, kicker=None):
    text(s, 0.7, 1.6, 8.6, 0.6, [{"runs": [(t, 25, True, NAVY, HEAD)], "sa": 0}])
    if kicker:
        text(s, 9.1, 1.76, 3.53, 0.35,
             [{"runs": [(kicker, 12, True, BLUE, HEAD)], "align": PP_ALIGN.RIGHT, "sa": 0}])


def footer(s, t, chip="concept"):
    # sit on the navy wave at the bottom-left so white text stays legible
    text(s, 0.7, 7.13, 8.6, 0.32,
         [{"runs": [(t, 9.5, False, WHITE, BODY)], "sa": 0}])
    if chip:
        mode_chip(s, chip)


def code_box(s, x, y, w, h, lines, sz=11):
    rrect(s, x, y, w, h, CODEBG, radius=0.06)
    paras = [{"runs": [(ln, sz, False, CODETX, CODE)], "sa": 2, "lh": 1.0} for ln in lines]
    text(s, x + 0.18, y + 0.14, w - 0.36, h - 0.28, paras)


def arrow(s, x, y, w=0.5, color=BLUE):
    text(s, x, y, w, 0.4, [{"runs": [("→", 24, True, color, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}],
         anchor=MSO_ANCHOR.MIDDLE)


def mode_chip(s, kind="concept"):
    """bottom-right pill: are learners absorbing (concept) or doing (lab)?"""
    if kind == "lab":
        fill = RGBColor(0xEA, 0xF5, 0xEE); tc = GREEN; label = "\U0001f9ea  LAB"
    else:
        fill = RGBColor(0xE7, 0xEF, 0xFB); tc = BLUE;  label = "\U0001f4d6  CONCEPT"
    w, h = 1.62, 0.34
    x = SW - 0.7 - w; y = 7.11   # on the navy wave (right side), aligned with footer
    rrect(s, x, y, w, h, fill, radius=0.5)
    text(s, x, y, w, h, [{"runs": [(label, 10.5, True, tc, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}],
         anchor=MSO_ANCHOR.MIDDLE)


def link_pill(s, x, y, w, h, label, url=REPO_URL, fill=WHITE, line=BLUE, tc=BLUE, sz=11):
    """clickable pill — both the shape and its text open the URL."""
    shp = rrect(s, x, y, w, h, fill, line=line, lw=1.3, radius=0.5)
    shp.click_action.hyperlink.address = url
    tb = text(s, x, y, w, h, [{"runs": [(label, sz, True, tc, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}],
              anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.paragraphs[0].runs[0].hyperlink.address = url
    return shp


def lab_marker(n, title_th, open_what, steps, minutes, outcome, accent=GREEN, login=None, pull_url=None):
    """full-slide hands-on checkpoint: stop lecturing, open the notebook/tool."""
    s = slide()
    rrect(s, 0.9, 1.78, 0.18, 4.35, accent, radius=0.4)
    text(s, 1.4, 1.74, 7.5, 0.5, [{"runs": [("\U0001f9ea  LAB %d · ลงมือทำ" % n, 15, True, accent, HEAD)], "sa": 0}])
    rrect(s, 10.45, 1.7, 2.18, 0.5, RGBColor(0xEA, 0xF5, 0xEE), line=accent, lw=1.2, radius=0.5)
    text(s, 10.45, 1.7, 2.18, 0.5, [{"runs": [("⏱  %d นาที" % minutes, 13, True, accent, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.4, 2.26, 11.0, 0.9, [{"runs": [(title_th, 29, True, NAVY, HEAD)], "sa": 0}])
    if pull_url:
        # primary one-click action: open JupyterHub and auto-pull the repo (nbgitpuller)
        link_pill(s, 1.4, 3.12, 7.0, 0.6, "▶  คลิกเดียว — เปิด hub + ดึง repo อัตโนมัติ",
                  url=pull_url, fill=GREEN, line=GREEN, tc=WHITE, sz=13.5)
    else:
        text(s, 1.4, 3.22, 7.8, 0.5, [{"runs": [("เปิด: ", 15, True, INK, HEAD), (open_what, 15, True, BLUE, CODE)], "sa": 0}])
    link_pill(s, 9.5, 3.18, 3.13, 0.5, "🔗 GitHub repo", sz=12)
    step_w = 6.3 if login else 10.6
    for i, st in enumerate(steps):
        yy = 3.95 + i * 0.5
        circle(s, 1.45, yy, 0.26, accent)
        text(s, 1.45, yy, 0.26, 0.26, [{"runs": [(str(i + 1), 11, True, WHITE, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.86, yy - 0.02, step_w, 0.3, [{"runs": [(st, 14, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    if login:
        url, usr, pwd = login
        lx, ly, lw, lh = 8.5, 3.95, 4.13, 1.85
        rrect(s, lx, ly, lw, lh, NAVY, radius=0.08)
        text(s, lx + 0.3, ly + 0.16, lw - 0.6, 0.35, [{"runs": [("🔑 บัญชีกลาง — ใช้ร่วมกัน", 12.5, True, RGBColor(0x9F, 0xC0, 0xEF), HEAD)], "sa": 0}])
        for j, (k, v) in enumerate([("URL", url), ("user", usr), ("pass", pwd)]):
            cy = ly + 0.62 + j * 0.4
            text(s, lx + 0.3, cy, 0.7, 0.32, [{"runs": [(k, 11.5, True, RGBColor(0xC9, 0xDA, 0xF3), HEAD)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
            text(s, lx + 0.95, cy, lw - 1.18, 0.32, [{"runs": [(v, 11.5, True, WHITE, CODE)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    oy = 3.95 + len(steps) * 0.5 + 0.2
    rrect(s, 1.4, oy, 11.23, 0.66, RGBColor(0xEF, 0xF7, 0xF1), line=accent, lw=1, radius=0.12)
    text(s, 1.7, oy, 10.65, 0.66, [{"runs": [("✅ เสร็จแล้วได้: ", 13.5, True, accent, HEAD), (outcome, 13, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Day 3 • ลงมือทำ — LAB %d" % n, chip=None)
    return s


# ============================================================ 1 TITLE
s = slide(COVER)
text(s, 0.9, 2.5, 6.7, 0.45, [{"runs": [("DAY 3  •  WORKSHOP", 14, True, CYANL, HEAD)], "sa": 0}])
text(s, 0.9, 3.0, 6.9, 1.4, [
    {"runs": [("Advanced Data Engineering", 34, True, WHITE, HEAD)], "sa": 4},
    {"runs": [("& Applied Analytics", 34, True, WHITE, HEAD)], "sa": 0},
])
text(s, 0.9, 4.5, 6.0, 0.5, [{"runs": [("Loan Analytics Pipeline", 16, False, LIGHT, BODY)], "sa": 0}])
rrect(s, 0.9, 5.12, 6.0, 0.6, None, line=CYANL, lw=1.5, radius=0.5)
text(s, 0.9, 5.12, 6.0, 0.6, [{"runs": [("CSV → Python ETL → MSSQL → Power BI", 13.5, True, CYANL, CODE)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 2 AGENDA
s = slide()
title(s, "โครงวันนี้ — 1 วัน (6 ชม.)", "AGENDA")
# morning card
rrect(s, 0.7, 2.0, 5.85, 3.7, CARD, line=CARDLN, lw=1)
rect(s, 0.7, 2.0, 5.85, 0.12, ORANGE)
text(s, 1.0, 2.28, 5.3, 0.5, [{"runs": [("\U0001f305  เช้า (3 ชม.)", 19, True, NAVY, HEAD)], "sa": 0}])
text(s, 1.0, 2.78, 5.3, 0.45, [{"runs": [("ETL → MSSQL Star Schema", 14, True, ORANGE, HEAD)], "sa": 0}])
for i, t in enumerate(["สำรวจข้อมูล + ออกแบบ star schema",
                       "ETL: clean → validate → transform",
                       "สร้าง dimensions + fact",
                       "โหลดเข้า MSSQL (de_loan_dw)"]):
    circle(s, 1.0, 3.42 + i * 0.52, 0.12, ORANGE)
    text(s, 1.25, 3.34 + i * 0.52, 5.0, 0.4, [{"runs": [(t, 13.5, False, INK, BODY)], "sa": 0}])
# afternoon card
rrect(s, 6.78, 2.0, 5.85, 3.7, CARD, line=CARDLN, lw=1)
rect(s, 6.78, 2.0, 5.85, 0.12, BLUE)
text(s, 7.08, 2.28, 5.3, 0.5, [{"runs": [("\U0001f306  บ่าย (3 ชม.)", 19, True, NAVY, HEAD)], "sa": 0}])
text(s, 7.08, 2.78, 5.3, 0.45, [{"runs": [("Power BI Dashboard", 14, True, BLUE, HEAD)], "sa": 0}])
for i, t in enumerate(["ต่อ Power BI → MSSQL",
                       "สร้าง model (star) + relationships",
                       "DAX measures: Default Rate %, ROI %",
                       "Dashboard: Analytics + Data Quality"]):
    circle(s, 7.08, 3.42 + i * 0.52, 0.12, BLUE)
    text(s, 7.33, 3.34 + i * 0.52, 5.0, 0.4, [{"runs": [(t, 13.5, False, INK, BODY)], "sa": 0}])
footer(s, "Day 3 • Agenda", chip=None)

# ============================================================ 3 BUSINESS PROBLEM
s = slide()
title(s, "โจทย์: Lending Club Loan Data", "บริบท")
text(s, 0.7, 2.05, 7.5, 1.05, [
    {"runs": [("Lending Club", 15, True, NAVY, HEAD), (" = แพลตฟอร์มสินเชื่อบุคคล (P2P) รายใหญ่ของสหรัฐฯ", 14, False, INK, BODY)], "sa": 3},
    {"runs": [("เปิดข้อมูลสินเชื่อย้อนหลังจริง — แต่ละแถว = 1 สินเชื่อ (วงเงิน · ดอกเบี้ย · เกรด A–G · สถานะ)", 12.5, False, MUTED, BODY)], "sa": 3},
    {"runs": [("เราใช้จำลองงาน DE จริง: เปลี่ยน ", 13, False, INK, BODY), ("ข้อมูลดิบ → ตอบคำถามธุรกิจ", 13, True, BLUE, HEAD), (" ได้", 13, False, INK, BODY)], "sa": 0}])
text(s, 0.7, 3.18, 7.4, 0.4, [{"runs": [("คำถามธุรกิจที่ dashboard จะตอบ:", 15, True, NAVY, HEAD)], "sa": 0}])
qs = [("\U0001f534", "เกรดไหน ", "เบี้ยวหนี้สูง", " เมื่อเทียบกับดอกเบี้ยที่คิด? → ตั้งราคาถูกไหม"),
      ("\U0001f4b0", "สินเชื่อกลุ่มไหน ", "กำไร/ขาดทุน", "?"),
      ("\U0001f4c8", "แนวโน้ม ", "ยอดปล่อยกู้ + อัตราเบี้ยวหนี้", " เป็นยังไง?")]
for i, (em, a, b, c) in enumerate(qs):
    yy = 3.6 + i * 0.66
    rrect(s, 0.7, yy, 7.4, 0.58, WHITE, line=CARDLN, lw=1, radius=0.16)
    text(s, 0.92, yy, 0.5, 0.58, [{"runs": [(em, 14, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.45, yy, 6.5, 0.58, [{"runs": [(a, 13, False, INK, BODY), (b, 13, True, NAVY, HEAD), (c, 13, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
# stat panel right
rrect(s, 8.45, 1.95, 4.18, 3.9, NAVY, radius=0.06)
text(s, 8.75, 2.35, 3.6, 0.9, [
    {"runs": [("ข้อมูลตัวอย่าง", 13, True, RGBColor(0x9F,0xC0,0xEF), HEAD)], "sa": 2},
    {"runs": [("loan_sample.csv", 15, True, WHITE, CODE)], "sa": 0}])
text(s, 8.75, 3.35, 3.6, 1.0, [
    {"runs": [("~50,000", 40, True, WHITE, HEAD)], "sa": 0},
    {"runs": [("แถว × ~25 คอลัมน์", 13, False, RGBColor(0xC9,0xDA,0xF3), BODY)], "sa": 0}])
text(s, 8.75, 4.7, 3.6, 0.9, [
    {"runs": [("ตัวอย่างคอลัมน์ที่ต้องล้าง:", 12, True, RGBColor(0x9F,0xC0,0xEF), HEAD)], "sa": 3},
    {"runs": [("int_rate \"13.56%\" · term \" 36 months\"", 11.5, False, RGBColor(0xDD,0xE8,0xF8), CODE)], "sa": 0}])
footer(s, "Day 3 • รู้จักข้อมูลก่อนลงมือ")

# ============================================================ 4 ARCHITECTURE
s = slide()
title(s, "ภาพรวม: จากไฟล์ดิบ→ dashboard", "ARCHITECTURE")
flow = [("CSV", "ข้อมูลดิบ (loan)", MUTED),
        ("Python ETL", "clean + transform", BLUE),
        ("Star Schema", "fact + 7 dims", CYAN),
        ("MSSQL", "de_loan_dw", NAVY),
        ("Power BI", "dashboard", ORANGE)]
bx, by, bw, bh = 0.72, 3.0, 2.18, 1.5
gap = 0.36
for i, (t, sub, col) in enumerate(flow):
    x = bx + i * (bw + gap)
    rrect(s, x, by, bw, bh, WHITE, line=col, lw=2, radius=0.1)
    rect(s, x, by, bw, 0.14, col)
    text(s, x + 0.1, by + 0.42, bw - 0.2, 0.5, [{"runs": [(t, 15, True, NAVY, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}])
    text(s, x + 0.1, by + 0.92, bw - 0.2, 0.4, [{"runs": [(sub, 11, False, MUTED, BODY)], "align": PP_ALIGN.CENTER, "sa": 0}])
    if i < 4:
        arrow(s, x + bw, by + 0.55, gap, BLUE)
text(s, 0.72, 2.1, 11.9, 0.5, [{"runs": [("ทุกอย่างรันบน ", 15, False, INK, BODY), ("JupyterHub", 15, True, NAVY, HEAD), (" → โหลดเข้า ", 15, False, INK, BODY), ("MSSQL (sql1)", 15, True, NAVY, HEAD), (" → Power BI ต่อผ่าน host", 15, False, INK, BODY)], "sa": 0}])
text(s, 0.72, 5.0, 11.9, 0.5, [{"runs": [("\U0001f4a1 โค้ด ETL = แพ็กเกจ ", 13, False, INK, BODY), ("loan_etl", 13, True, BLUE, CODE), (" — ชุดเดียวใช้ได้ทั้ง notebook / CLI / API", 13, False, INK, BODY)], "sa": 0}])
# define the ETL acronym (tested in Q1)
rrect(s, 0.72, 5.55, 11.9, 0.62, RGBColor(0xEE, 0xF3, 0xFB), line=CARDLN, lw=1, radius=0.1)
text(s, 1.02, 5.55, 11.4, 0.62, [{"runs": [("ETL = ", 13.5, True, NAVY, HEAD), ("Extract", 13, True, BLUE, HEAD), (" ดึงข้อมูลออกมา → ", 12.5, False, INK, BODY), ("Transform", 13, True, BLUE, HEAD), (" แปลง/ล้าง → ", 12.5, False, INK, BODY), ("Load", 13, True, BLUE, HEAD), (" โหลดเข้าปลายทาง", 12.5, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Day 3 • End-to-end pipeline")

# ============================================================ 5 DIVIDER morning
def divider(title_th, sub, accent):
    s = slide()
    rrect(s, 0.9, 2.7, 0.18, 2.1, accent, radius=0.4)
    text(s, 1.35, 2.7, 10.8, 0.6, [{"runs": [("\U0001f305  MORNING SESSION", 15, True, accent, HEAD)], "sa": 0}])
    text(s, 1.35, 3.25, 11.0, 1.2, [{"runs": [(title_th, 36, True, NAVY, HEAD)], "sa": 0}])
    text(s, 1.35, 4.45, 11.0, 0.5, [{"runs": [(sub, 16, False, INK, BODY)], "sa": 0}])
    return s
s = divider("เช้า — ETL → MSSQL Star Schema", "จาก CSV ดิบ → ข้อมูลสะอาดใน data warehouse", ORANGE)

# ============================================================ 6 SETUP
s = slide()
title(s, "เครื่องมือ + Setup", "เช้า · ขั้นเตรียม")
cards = [("\U0001f4bb", "JupyterHub", "รัน notebook ทำ ETL บน browser — ไม่ต้องติดตั้งอะไรในเครื่อง", BLUE),
         ("\U0001f4e6", "loan_etl", "แพ็กเกจ ETL พร้อมใช้ — clean / validate / transform / load", CYAN),
         ("\U0001f517", "nbgitpuller", "คลิกลิงก์เดียว ดึง repo + notebook มาไว้บน hub", NAVY)]
for i, (em, h, d, col) in enumerate(cards):
    x = 0.7 + i * 4.07
    rrect(s, x, 2.2, 3.85, 2.5, CARD, line=CARDLN, lw=1)
    circle(s, x + 0.3, 2.5, 0.62, col)
    text(s, x + 0.3, 2.48, 0.62, 0.62, [{"runs": [(em, 18, False, WHITE, BODY)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 1.1, 2.62, 2.6, 0.5, [{"runs": [(h, 17, True, NAVY, HEAD)], "sa": 0}])
    text(s, x + 0.3, 3.4, 3.25, 1.2, [{"runs": [(d, 12.5, False, INK, BODY)], "sa": 0, "lh": 1.1}])
rrect(s, 0.7, 5.0, 8.3, 0.92, NAVY, radius=0.1)
text(s, 1.05, 5.13, 7.6, 0.32, [{"runs": [("🔑 เข้าใช้ JupyterHub — ทุกคนใช้บัญชีกลางนี้", 13, True, RGBColor(0x9F, 0xC0, 0xEF), HEAD)], "sa": 0}])
text(s, 1.05, 5.5, 7.6, 0.35, [{"runs": [
    ("URL ", 12, True, RGBColor(0xC9, 0xDA, 0xF3), HEAD), ("jupyter-hub.minddatatech.com", 12, True, WHITE, CODE),
    ("   user ", 12, True, RGBColor(0xC9, 0xDA, 0xF3), HEAD), ("deadmin", 12, True, WHITE, CODE),
    ("   pass ", 12, True, RGBColor(0xC9, 0xDA, 0xF3), HEAD), ("deadmin@2026", 12, True, WHITE, CODE)], "sa": 0}])
# one-click nbgitpuller button (same action as LAB 1) for early access
link_pill(s, 9.15, 5.0, 3.48, 0.92, "▶  ดึง repo เข้า hub", url=PULL_URL, fill=GREEN, line=GREEN, tc=WHITE, sz=14)
footer(s, "Day 3 • เช้า — Setup")

# ============================================================ 7 EXPLORE
s = slide()
title(s, "สำรวจข้อมูล — ของจริงไม่เคยสะอาด", "Notebook 01")
text(s, 0.7, 2.1, 11.9, 0.4, [{"runs": [("คอลัมน์เก็บมาเป็น “ข้อความปนสัญลักษณ์” — คำนวณไม่ได้ทันที:", 14, True, NAVY, HEAD)], "sa": 0}])
rows = [("int_rate", '"13.56%"', "13.56 (ตัวเลข)"),
        ("term", '" 36 months"', "36"),
        ("emp_length", '"10+ years"', "10"),
        ("issue_d", '"Dec-2018"', "วันที่จริง")]
ry = 2.55
rrect(s, 0.7, ry, 7.3, 0.5, NAVY, radius=0.05)
for cx, w, txt2 in [(0.95, 2.3, "คอลัมน์"), (3.3, 2.6, "ดิบ (ข้อความ)"), (6.0, 1.9, "หลังแก้")]:
    text(s, cx, ry, w, 0.5, [{"runs": [(txt2, 12.5, True, WHITE, HEAD)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
for i, (c, raw, fixed) in enumerate(rows):
    yy = ry + 0.5 + i * 0.52
    fill = WHITE if i % 2 == 0 else CARD
    rect(s, 0.7, yy, 7.3, 0.52, fill)
    text(s, 0.95, yy, 2.3, 0.52, [{"runs": [(c, 12, True, BLUE, CODE)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.3, yy, 2.6, 0.52, [{"runs": [(raw, 12, False, ORANGE, CODE)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.0, yy, 1.9, 0.52, [{"runs": [(fixed, 12, False, GREEN, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
rrect(s, 0.7, ry, 7.3, 0.5 + 4 * 0.52, None, line=CARDLN, lw=1, radius=0.05)
# goal panel
rrect(s, 8.3, 2.55, 4.33, 2.6, RGBColor(0xEE,0xF3,0xFB), line=BLUE, lw=1.2, radius=0.08)
text(s, 8.6, 2.8, 3.8, 0.5, [{"runs": [("\U0001f3af เป้าหมาย", 16, True, NAVY, HEAD)], "sa": 0}])
text(s, 8.6, 3.3, 3.8, 1.7, [
    {"runs": [("ปั้นข้อมูลให้เป็น ", 13.5, False, INK, BODY), ("Star Schema", 13.5, True, BLUE, HEAD)], "sa": 6},
    {"runs": [("• 1 fact (fact_loan)", 13, False, INK, BODY)], "sa": 3},
    {"runs": [("• 7 dimensions (date, grade,", 13, False, INK, BODY)], "sa": 1},
    {"runs": [("   purpose, status, geo, …)", 13, False, INK, BODY)], "sa": 0}])
footer(s, "Day 3 • เช้า — Explore")

# ============================================================ LAB 1
lab_marker(1, "ตั้งสภาพแวดล้อม + สำรวจข้อมูลจริง", "JupyterHub → notebooks/00, 01",
           ["login ด้วยบัญชีกลาง (กล่องขวา →)",
            "กดปุ่มเขียวด้านบน — ดึง repo อัตโนมัติ",
            "รัน 00_setup — ตั้ง TEAM_ID ของทีม",
            "รัน 01_explore — ดูข้อมูลดิบที่ต้องล้าง"],
           20, "repo อยู่บน hub · เห็นปัญหาในข้อมูล · พร้อมลงมือ ETL",
           login=("jupyter-hub.minddatatech.com", "deadmin", "deadmin@2026"),
           pull_url=PULL_URL)

# ============================================================ DATA MODEL (why star schema)
s = slide()
title(s, "Data Model — ทำไมต้อง Star Schema?", "เช้า · ออกแบบก่อนทำ")
text(s, 0.7, 2.1, 11.9, 0.5, [{"runs": [("ก่อนเขียน ETL ออกแบบ ", 15, False, INK, BODY), ("“ปลายทาง”", 15, True, NAVY, HEAD), (" ก่อน — ข้อมูลควรอยู่รูปไหนถึงตอบคำถามธุรกิจได้ง่ายที่สุด", 15, False, INK, BODY)], "sa": 0}])
_con = [("fact (ตารางกลาง)", ["เก็บ “ตัวเลขวัดผล” — วงเงิน · ดอกเบี้ย", "is_default · profit", "1 แถว = 1 สินเชื่อ + คีย์ไปหา dim"], NAVY),
        ("dimension (ตารางรอบ)", ["เก็บ “มุมมอง” — เกรด · วันที่ · รัฐ", "วัตถุประสงค์", "ใช้ filter / group เวลาวิเคราะห์"], BLUE)]
for i, (h, lines, col) in enumerate(_con):
    x = 0.7 + i * 6.1
    rrect(s, x, 2.82, 5.83, 1.78, CARD, line=CARDLN, lw=1)
    rect(s, x, 2.82, 0.14, 1.78, col)
    text(s, x + 0.38, 3.0, 5.3, 0.45, [{"runs": [(h, 16, True, NAVY, HEAD)], "sa": 0}])
    paras = [{"runs": [(ln, 12, False, INK, BODY)], "sa": 3, "lh": 1.05} for ln in lines]
    text(s, x + 0.38, 3.55, 5.2, 1.0, paras)
rrect(s, 0.7, 4.82, 11.93, 1.12, RGBColor(0xEE, 0xF3, 0xFB), line=BLUE, lw=1, radius=0.1)
text(s, 1.0, 4.95, 11.4, 0.4, [{"runs": [("ทำไม star schema?", 14, True, BLUE, HEAD)], "sa": 0}])
text(s, 1.0, 5.32, 11.4, 0.55, [
    {"runs": [("• ตอบคำถามแบบ ", 12.5, False, INK, BODY), ("“X ต่อ Y”", 12.5, True, NAVY, HEAD), (" ได้ทันที — default rate ต่อเกรด · กำไรต่อวัตถุประสงค์", 12.5, False, INK, BODY)], "sa": 2},
    {"runs": [("• Power BI ทำงานกับ star schema ได้เร็ว + relationship เข้าใจง่าย", 12.5, False, INK, BODY)], "sa": 0}])
footer(s, "Day 3 • Data Model")

# ============================================================ STAR SCHEMA (diagram)
s = slide()
title(s, "Star Schema — fact + 7 dimensions", "เช้า · Data Model")
fx, fy, fw, fh = 5.35, 3.5, 2.65, 1.35
DW, DH = 1.83, 0.52
dims_left = [("dim_date", 0.85, 2.35), ("dim_purpose", 0.85, 3.85), ("dim_borrower", 0.85, 5.35)]
dims_right = [("dim_loan_status", 9.65, 2.35), ("dim_geography", 9.65, 3.85), ("dim_term", 9.65, 5.35)]
dim_top = ("dim_grade", 5.755, 2.12)
SPOKE = RGBColor(0x7F, 0xA8, 0xDE)
def spoke(x1, y1, x2, y2):
    ln = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = SPOKE; ln.line.width = Pt(1.75); ln.shadow.inherit = False
edge_y = [fy + 0.32, fy + fh / 2, fy + fh - 0.32]
for (nm, dx, dy), py in zip(dims_left, edge_y):
    spoke(dx + DW, dy + DH / 2, fx, py)
for (nm, dx, dy), py in zip(dims_right, edge_y):
    spoke(dx, dy + DH / 2, fx + fw, py)
spoke(dim_top[1] + DW / 2, dim_top[2] + DH, fx + fw / 2, fy)
rrect(s, fx - 0.09, fy - 0.09, fw + 0.18, fh + 0.18, RGBColor(0xDC, 0xE8, 0xF8), radius=0.16)
rrect(s, fx, fy, fw, fh, NAVY, line=CYAN, lw=1.5, radius=0.12)
text(s, fx, fy + 0.24, fw, 0.5, [{"runs": [("★ fact_loan", 18, True, WHITE, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}])
text(s, fx, fy + 0.76, fw, 0.5, [{"runs": [("loan_amnt · int_rate", 10.5, False, RGBColor(0xCF, 0xDF, 0xF6), CODE)], "align": PP_ALIGN.CENTER, "sa": 0}, {"runs": [("is_default · profit …", 10.5, False, RGBColor(0xCF, 0xDF, 0xF6), CODE)], "align": PP_ALIGN.CENTER, "sa": 0}])
for nm, dx, dy in dims_left + dims_right + [dim_top]:
    rrect(s, dx, dy, DW, DH, WHITE, line=BLUE, lw=1.25, radius=0.16)
    text(s, dx, dy, DW, DH, [{"runs": [(nm, 11.5, True, NAVY, CODE)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
# Grain explained as a concept (tested in Q5) — was just a faint caption before
rrect(s, 1.7, 6.16, 9.93, 0.66, RGBColor(0xEE, 0xF3, 0xFB), line=BLUE, lw=1, radius=0.14)
text(s, 2.0, 6.16, 9.33, 0.66, [{"runs": [("📐 Grain = ", 13, True, BLUE, HEAD), ("“1 แถวของ fact แทนอะไร”", 12.5, True, NAVY, HEAD), (" → ที่นี่ = 1 สินเชื่อ  ·  กำหนดให้ชัดเป็นขั้นแรกของการออกแบบ schema", 12, False, INK, BODY)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
mode_chip(s)

# ============================================================ 8 CLEAN
s = slide()
title(s, "ETL ① Clean — แก้ format ให้คำนวณได้", "Notebook 02")
text(s, 0.7, 2.12, 6.3, 1.4, [
    {"runs": [("เปลี่ยนข้อความปนสัญลักษณ์ → ตัวเลข/วันที่จริง", 15, False, INK, BODY)], "sa": 8},
    {"runs": [("•  ", 14, True, BLUE, BODY), ("ลบ % / ตัดคำว่า months / years", 13.5, False, INK, BODY)], "sa": 5},
    {"runs": [("•  ", 14, True, BLUE, BODY), ("แปลงเป็น numeric / datetime", 13.5, False, INK, BODY)], "sa": 5},
    {"runs": [("•  ", 14, True, BLUE, BODY), ("ตัดช่องว่าง (strip) ข้อความ", 13.5, False, INK, BODY)], "sa": 0}])
code_box(s, 7.2, 2.12, 5.43, 3.5, [
    "# int_rate \"13.56%\" -> 13.56",
    "df['int_rate'] = pd.to_numeric(",
    "    df['int_rate'].str",
    "      .replace('%','').str.strip())",
    "",
    "# term \" 36 months\" -> 36",
    "df['term_months'] = (df['term']",
    "    .str.extract(r'(\\d+)'))",
    "",
    "# issue_d \"Dec-2018\" -> datetime",
    "df['issue_date'] = pd.to_datetime(",
    "    df['issue_d'], format='%b-%Y')",
], sz=11.5)
footer(s, "Day 3 • ETL — Clean")

# ============================================================ 9 VALIDATE
s = slide()
title(s, "ETL ② Validate + Quarantine", "Notebook 02")
text(s, 0.7, 2.12, 11.9, 0.45, [{"runs": [("ตรวจ ", 15, False, INK, BODY), ("ทุกแถว", 15, True, NAVY, HEAD), (" ตามกฎ business — แถวไม่ผ่าน “ไม่ลบ” แต่เก็บใน quarantine พร้อมเหตุผล (audit)", 14, False, INK, BODY)], "sa": 0}])
code_box(s, 0.7, 2.68, 6.45, 2.92, [
    "# รวมทุกกฎเป็น mask เดียว",
    "mask = ( (df['loan_amnt'] > 0)",
    "  & (df['annual_inc'] >= 0)",
    "  & df['term_months'].isin([36,60])",
    "  & df['loan_status'].isin(VALID) )",
    "",
    "good = df[mask]      # ผ่าน",
    "bad  = df[~mask]     # ไม่ผ่าน",
    "bad['reject_reason'] = why(bad)",
], sz=11)
text(s, 7.4, 2.68, 5.25, 0.4, [{"runs": [("แต่ละบรรทัดทำอะไร", 14, True, NAVY, HEAD)], "sa": 0}])
_ve = [("mask", "เงื่อนไขทุกข้อรวมด้วย & (AND)"),
       ("df[mask]", "เก็บเฉพาะแถวที่ผ่าน → good"),
       ("df[~mask]", "แถวไม่ผ่าน (~ = NOT) → bad"),
       ("reject_reason", "ใส่เหตุผลที่ถูกคัด → audit")]
for i, (a, b) in enumerate(_ve):
    text(s, 7.4, 3.2 + i * 0.55, 5.25, 0.5, [{"runs": [(a, 12, True, BLUE, CODE), ("  " + b, 12, False, INK, BODY)], "sa": 0, "lh": 1.0}])
rrect(s, 7.4, 5.42, 5.25, 0.5, RGBColor(0xEF, 0xF7, 0xF1), line=GREEN, lw=1, radius=0.12)
text(s, 7.4, 5.42, 5.25, 0.5, [{"runs": [("ผล: good 482  ·  quarantine 18 แถว", 12.5, True, NAVY, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Day 3 • ETL — Validate")

# ============================================================ 10 TRANSFORM
s = slide()
title(s, "ETL ③ Transform — สร้างความหมายธุรกิจ", "Notebook 02")
text(s, 0.7, 2.12, 11.9, 0.45, [{"runs": [("เพิ่มคอลัมน์ที่ ", 15, False, INK, BODY), ("ธุรกิจใช้จริง", 15, True, NAVY, HEAD), (" จากข้อมูลที่สะอาดแล้ว — ขั้นนี้คือ ~70% ของงาน DE", 14, False, INK, BODY)], "sa": 0}])
code_box(s, 0.7, 2.68, 6.45, 2.95, [
    "def derive(df):",
    "  df['is_default'] = (df['loan_status']",
    "      .isin(BAD).astype(int))",
    "  df['fico_avg'] = df[['fico_low',",
    "      'fico_high']].mean(axis=1)",
    "  df['profit'] = (df['total_pymnt']",
    "      - df['funded_amnt'])",
    "  return df",
], sz=11)
text(s, 7.4, 2.68, 5.25, 0.4, [{"runs": [("คอลัมน์ที่ได้ (ใช้ทำ KPI)", 14, True, NAVY, HEAD)], "sa": 0}])
_te = [("is_default", "1 = เบี้ยวหนี้ → คำนวณ default rate", ORANGE),
       ("fico_avg", "คะแนนเครดิตเฉลี่ย (low+high)/2", CYAN),
       ("profit", "กำไร = เงินได้คืน − เงินปล่อยกู้", GREEN),
       ("status_group", "จัดกลุ่มสถานะ Good / Bad", PURPLE)]
for i, (c, d, col) in enumerate(_te):
    yy = 3.25 + i * 0.62
    circle(s, 7.42, yy + 0.06, 0.13, col)
    text(s, 7.72, yy, 5.0, 0.55, [{"runs": [(c, 12, True, NAVY, CODE), ("  " + d, 11.5, False, INK, BODY)], "sa": 0, "lh": 1.0}])
footer(s, "Day 3 • ETL — Transform")

# (Star Schema diagram moved up to the Data Model section, before ETL)

# ============================================================ LAB 2
lab_marker(2, "ETL: clean → validate → transform", "notebooks/02_etl",
           ["รัน clean — แปลง %, term, วันที่ ให้คำนวณได้",
            "รัน validate — แยก good / quarantine + เหตุผล",
            "รัน transform — สร้าง is_default · profit · fico_avg",
            "ตรวจจำนวนแถว good / quarantine ว่าตรงไหม"],
           30, "DataFrame สะอาด + คอลัมน์ธุรกิจ พร้อมขึ้น star schema")

# ============================================================ 12 LOAD MSSQL
s = slide()
title(s, "โหลดเข้า MSSQL — จบ ETL", "Notebook 03")
items = [("\U0001f5c4️", "Database", "de_loan_dw", BLUE),
         ("\U0001f464", "Login (ร่วมกัน)", "deadmin", NAVY),
         ("\U0001f3f7️", "ชื่อตาราง", "<TEAM_ID>_fact_loan", CYAN),
         ("✅", "Reconcile", "row + sum ตรงกัน", GREEN)]
for i, (em, h, v, col) in enumerate(items):
    x = 0.7 + (i % 2) * 6.1
    y = 2.15 + (i // 2) * 1.45
    rrect(s, x, y, 5.83, 1.25, WHITE, line=CARDLN, lw=1)
    circle(s, x + 0.28, y + 0.32, 0.6, col)
    text(s, x + 0.28, y + 0.3, 0.6, 0.6, [{"runs": [(em, 17, False, WHITE, BODY)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 1.05, y + 0.22, 4.6, 0.4, [{"runs": [(h, 13, True, MUTED, HEAD)], "sa": 0}])
    text(s, x + 1.05, y + 0.6, 4.6, 0.5, [{"runs": [(v, 17, True, NAVY, CODE)], "sa": 0}])
rrect(s, 0.7, 5.1, 11.93, 0.85, RGBColor(0xEF,0xF7,0xF1), line=GREEN, lw=1, radius=0.12)
text(s, 1.0, 5.1, 11.4, 0.85, [{"runs": [("✅ ทุกคนใช้ ", 14, True, GREEN, HEAD), ("deadmin", 13.5, True, NAVY, CODE), (" เหมือนกัน — กันชนด้วย prefix ชื่อตาราง (TEAM_ID) → ไม่ต้องสร้าง schema รายทีม", 13.5, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Day 3 • โหลดเข้า MSSQL")

# ============================================================ LAB 3
lab_marker(3, "โหลด Star Schema เข้า MSSQL", "notebooks/03_star_schema_load",
           ["สร้าง dimensions + fact จากข้อมูลที่ transform แล้ว",
            "เขียนลง de_loan_dw → ตาราง <TEAM_ID>_*",
            "รัน reconcile — row count + ยอดรวมต้องตรง",
            "เปิด Power BI Desktop ค้างไว้ ใช้ต่อช่วงบ่าย"],
           25, "fact_loan + 7 dimensions อยู่ใน MSSQL — จบกระบวนการ ETL")

# ============================================================ 13 DIVIDER afternoon
s = slide()
rrect(s, 0.9, 2.7, 0.18, 2.1, BLUE, radius=0.4)
text(s, 1.35, 2.7, 10.8, 0.6, [{"runs": [("\U0001f306  AFTERNOON SESSION", 15, True, BLUE, HEAD)], "sa": 0}])
text(s, 1.35, 3.25, 11.0, 1.2, [{"runs": [("บ่าย — Power BI Dashboard", 36, True, NAVY, HEAD)], "sa": 0}])
text(s, 1.35, 4.45, 11.0, 0.5, [{"runs": [("ต่อ MSSQL → สร้าง model + DAX → dashboard เล่าเรื่อง", 16, False, INK, BODY)], "sa": 0}])

# ============================================================ INSTALL POWER BI
s = slide()
title(s, "ติดตั้ง Power BI Desktop (Windows)", "บ่าย · เตรียม")
text(s, 0.7, 2.1, 11.9, 0.4, [{"runs": [("ฟรี · ใช้ได้เฉพาะ Windows — เลือกติดตั้งได้ 2 ทาง", 14, False, INK, BODY)], "sa": 0}])
_inst = [("1", "Microsoft Store", ["เปิด Microsoft Store", "ค้นหา \"Power BI Desktop\"", "กด Get / Install", "อัปเดตอัตโนมัติ"], BLUE, "แนะนำ"),
         ("2", "ดาวน์โหลด Installer", ["powerbi.microsoft.com/desktop", "Download → ได้ไฟล์ .exe", "รันไฟล์ → Next → Install", "เหมาะกับเครื่องที่ปิด Store"], CYAN, "")]
for i, (n, h, steps, col, tag) in enumerate(_inst):
    x = 0.7 + i * 6.1
    rrect(s, x, 2.65, 5.83, 2.72, CARD, line=CARDLN, lw=1)
    badge(s, x + 0.3, 2.9, 0.6, n, col, WHITE, 19)
    text(s, x + 1.05, 3.0, 3.3, 0.45, [{"runs": [(h, 16, True, NAVY, HEAD)], "sa": 0}])
    if tag:
        rrect(s, x + 4.55, 2.96, 1.0, 0.4, GREEN, radius=0.5)
        text(s, x + 4.55, 2.96, 1.0, 0.4, [{"runs": [(tag, 11, True, WHITE, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    for j, st in enumerate(steps):
        circle(s, x + 0.35, 3.76 + j * 0.37, 0.1, col)
        text(s, x + 0.6, 3.68 + j * 0.37, 5.0, 0.35, [{"runs": [(st, 12, False, INK, BODY)], "sa": 0}])
rrect(s, 0.7, 5.56, 11.93, 0.68, RGBColor(0xFC, 0xF1, 0xE7), line=ORANGE, lw=1, radius=0.12)
text(s, 1.0, 5.56, 11.4, 0.68, [{"runs": [("⚠️ ", 13, False, ORANGE, BODY), ("Mac / Linux: ", 13, True, NAVY, HEAD), ("ไม่มี Power BI Desktop — ใช้ ", 12.5, False, INK, BODY), ("Power BI Service (เว็บ)", 12.5, True, NAVY, HEAD), (" หรือรันผ่าน Windows VM", 12.5, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Day 3 • บ่าย — ติดตั้ง Power BI")

# ============================================================ 14 POWER BI connect+model
s = slide()
title(s, "Power BI: Connect → Model → DAX", "บ่าย")
steps = [("1", "Connect", "Get Data → SQL Server\n34.59.223.196 / de_loan_dw\nเลือกตาราง <TEAM_ID>_*", BLUE),
         ("2", "Model", "เชื่อม relationships\nfact → dim (star)\nmark dim_date เป็น date table", CYAN),
         ("3", "DAX", "Default Rate %\nROI %, Total Funded\nAvg Int Rate", NAVY)]
for i, (n, h, body, col) in enumerate(steps):
    x = 0.7 + i * 4.07
    rrect(s, x, 2.15, 3.85, 3.4, CARD, line=CARDLN, lw=1)
    badge(s, x + 0.3, 2.42, 0.62, n, col, WHITE, 20)
    text(s, x + 1.1, 2.55, 2.6, 0.5, [{"runs": [(h, 18, True, NAVY, HEAD)], "sa": 0}])
    paras = [{"runs": [(ln, 12.5, False, INK, BODY)], "sa": 4, "lh": 1.05} for ln in body.split("\n")]
    text(s, x + 0.32, 3.35, 3.3, 2.0, paras)
# show the actual Default Rate % DAX formula (tested in Q10)
rrect(s, 0.7, 5.72, 11.93, 0.72, CODEBG, radius=0.08)
text(s, 1.0, 5.72, 11.4, 0.72, [{"runs": [("ตัวอย่าง DAX:   ", 12.5, True, CYANL, HEAD), ("Default Rate %", 12.5, True, WHITE, CODE), (" = ", 12.5, False, CODETX, CODE), ("DIVIDE( SUM(is_default), COUNTROWS(fact_loan) ) * 100", 12, False, CODETX, CODE)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Day 3 • บ่าย — Power BI model")

# ============================================================ LAB 4
lab_marker(4, "ติดตั้ง Power BI + ต่อ MSSQL", "Power BI Desktop → SQL Server",
           ["ติดตั้ง Power BI Desktop (Store หรือ .exe)",
            "Get Data → SQL Server → 34.59.223.196 / de_loan_dw",
            "เลือกตาราง <TEAM_ID>_* ของทีมเรา",
            "Model view: เชื่อม fact → dim + เพิ่ม DAX measures"],
           45, "model แบบ star + measures (Default Rate %, ROI %) พร้อมใช้")

# ============================================================ 15 DASHBOARD LAYOUT (mockup)
s = slide()
title(s, "Loan Analytics Dashboard — Layout", "บ่าย · ผลลัพธ์")
cx, cy, cw, ch = 0.7, 2.1, 11.93, 4.12
rrect(s, cx, cy, cw, ch, RGBColor(0xF3, 0xF7, 0xFC), line=CARDLN, lw=1, radius=0.05)
_kpi = ["Total Funded", "Total Loans", "Default Rate %", "Avg Int Rate", "ROI %"]
kw = (cw - 0.4) / 5
for i, k in enumerate(_kpi):
    kx = cx + 0.2 + i * kw
    rrect(s, kx, cy + 0.18, kw - 0.16, 0.82, WHITE, line=CARDLN, lw=1, radius=0.1)
    text(s, kx + 0.14, cy + 0.28, kw - 0.4, 0.3, [{"runs": [(k, 10, True, MUTED, HEAD)], "sa": 0}])
    text(s, kx + 0.14, cy + 0.56, kw - 0.4, 0.35, [{"runs": [("00.0", 18, True, NAVY, HEAD)], "sa": 0}])
_pan = [("Default Rate % vs Int Rate · by Grade", "bars", BLUE),
        ("Total Funded · by Month", "line", CYAN),
        ("Total Profit · by Purpose", "bars", GREEN),
        ("Default Rate % · by State", "map", ORANGE)]
pw = (cw - 0.6) / 2
ph = 1.34
for i, (t2, kind, col) in enumerate(_pan):
    px = cx + 0.2 + (i % 2) * (pw + 0.2)
    py = cy + 1.18 + (i // 2) * (ph + 0.12)
    rrect(s, px, py, pw, ph, WHITE, line=CARDLN, lw=1, radius=0.07)
    rect(s, px, py, 0.09, ph, col)
    text(s, px + 0.24, py + 0.12, pw - 0.4, 0.3, [{"runs": [(t2, 10.5, True, NAVY, HEAD)], "sa": 0}])
    if kind == "bars":
        for b in range(6):
            bh = 0.16 + 0.1 * ((b * 2 + 1) % 5 + 1)
            rect(s, px + 0.32 + b * 0.42, py + ph - 0.18 - bh, 0.27, bh, col)
    elif kind == "line":
        ys = [0.62, 0.42, 0.55, 0.30, 0.46, 0.22, 0.34]
        step = (pw - 0.7) / (len(ys) - 1)
        for a in range(len(ys) - 1):
            ln2 = s.shapes.add_connector(1, Inches(px + 0.32 + a * step), Inches(py + ph - 0.2 - ys[a]), Inches(px + 0.32 + (a + 1) * step), Inches(py + ph - 0.2 - ys[a + 1]))
            ln2.line.color.rgb = col; ln2.line.width = Pt(2.25); ln2.shadow.inherit = False
    elif kind == "map":
        circle(s, px + 0.55, py + 0.48, 0.62, RGBColor(0xF6, 0xE6, 0xD6))
        circle(s, px + 1.45, py + 0.62, 0.42, RGBColor(0xF1, 0xD8, 0xC0))
        text(s, px + 2.25, py + 0.5, pw - 2.45, 0.6, [{"runs": [("choropleth ตามรัฐ", 10.5, False, MUTED, BODY)], "sa": 0}])
footer(s, "Day 3 • Dashboard (Analytics)")

# ============================================================ DQ DASHBOARD — DAX (dq_power-bi)
s = slide()
title(s, "Data Quality Dashboard — DAX", "บ่าย · dq_power-bi")
text(s, 0.7, 2.1, 11.9, 0.4, [{"runs": [("วัดคุณภาพข้อมูลที่ ETL แล้ว 3 มุม — ", 14, False, INK, BODY), ("Data Rows · Validity · Completeness", 14, True, NAVY, HEAD)], "sa": 0}])
_m = [("Data Rows", "จำนวนแถวทั้งหมด", BLUE),
      ("Validity", "% แถวที่ผ่านกฎ (valid)", ORANGE),
      ("Completeness", "% ช่องที่ไม่ว่าง (not null)", GREEN)]
for i, (h, d, col) in enumerate(_m):
    yy = 2.72 + i * 0.92
    rrect(s, 0.7, yy, 5.4, 0.78, WHITE, line=CARDLN, lw=1, radius=0.12)
    rect(s, 0.7, yy, 0.14, 0.78, col)
    text(s, 1.05, yy + 0.12, 5.0, 0.4, [{"runs": [(h, 16, True, NAVY, HEAD)], "sa": 0}])
    text(s, 1.05, yy + 0.45, 5.0, 0.3, [{"runs": [(d, 12, False, MUTED, BODY)], "sa": 0}])
code_box(s, 6.4, 2.72, 6.23, 2.84, [
    "TotalRows = COUNTROWS(loan_data)",
    "",
    "ValidRows = CALCULATE(",
    "    COUNTROWS(loan_data),",
    "    loan_data[is_valid] = \"TRUE\")",
    "",
    "Validity %  = DIVIDE([ValidRows],",
    "                     [TotalRows]) * 100",
    "",
    "Completeness % = DIVIDE(",
    "    [NotNullCount], [TotalRows]) * 100",
], sz=10.5)
footer(s, "Day 3 • DQ Dashboard")

# ============================================================ DQ DASHBOARD — LAYOUT (dq_power-bi)
s = slide()
title(s, "DQ Dashboard — Layout & Visuals", "บ่าย · dq_power-bi")
text(s, 0.7, 2.1, 11.9, 0.4, [{"runs": [("3 ส่วนตาม dq_power-bi — ใช้ ", 14, False, INK, BODY), ("Card", 14, True, NAVY, HEAD), (" (ตัวเลขใหญ่) + ", 14, False, INK, BODY), ("Donut", 14, True, NAVY, HEAD), (" (สัดส่วน TRUE/FALSE/NULL)", 14, False, INK, BODY)], "sa": 0}])
_b = [("Data Rows", "Card แสดง TotalRows (เช่น 1.4M แถว)", DQ_BLUE, "#9DC3E6"),
      ("Validity", "Donut: TRUE / FALSE / NULL ต่อคอลัมน์", DQ_YEL, "#FFE699"),
      ("Completeness", "Donut: NULL vs NOT NULL ต่อคอลัมน์", DQ_GRN, "#C6E0B4")]
for i, (h, d, fill, hexc) in enumerate(_b):
    yy = 2.7 + i * 1.0
    rrect(s, 0.7, yy, 11.93, 0.84, fill, radius=0.1)
    text(s, 1.05, yy, 3.2, 0.84, [{"runs": [(h, 17, True, NAVY, HEAD)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.3, yy, 6.2, 0.84, [{"runs": [(d, 13, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 10.7, yy, 1.9, 0.84, [{"runs": [(hexc, 12, True, NAVY, CODE)], "align": PP_ALIGN.RIGHT, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 5.85, 11.9, 0.4, [{"runs": [("ต้นแบบเต็ม + ไฟล์ .pbix: ", 12.5, False, MUTED, BODY), ("github.com/amornpan/dq_power-bi", 12.5, True, BLUE, CODE)], "sa": 0}])
footer(s, "Day 3 • DQ Layout")

# ============================================================ LAB 5
lab_marker(5, "สร้าง Dashboard เล่าเรื่อง", "Power BI → Report view",
           ["หน้า Analytics: KPI cards + bar / line / map",
            "หน้า Data Quality: Card + Donut (validity / completeness)",
            "ใส่ slicer (เกรด · ช่วงเวลา) ให้กดเล่นได้",
            "ตรวจตัวเลขกับผล reconcile ให้ตรง"],
           60, "dashboard 2 หน้า ตอบโจทย์ธุรกิจ + วัดคุณภาพข้อมูลได้")

# ============================================================ 16 WRAP / THANKS
s = slide()
text(s, 1.0, 1.5, 11.3, 0.5, [{"runs": [("สิ่งที่คุณทำได้หลังวันนี้", 14, True, BLUE, HEAD)], "sa": 0}])
text(s, 1.0, 1.92, 11.3, 0.8, [{"runs": [("สร้าง pipeline จริง CSV → Power BI", 30, True, NAVY, HEAD)], "sa": 0}])
done = ["เขียน Python ETL: clean / validate / quarantine / transform",
        "ออกแบบ star schema → โหลดเข้า MSSQL (de_loan_dw)",
        "ต่อ Power BI → สร้าง dashboard (Analytics + Data Quality)"]
for i, t in enumerate(done):
    circle(s, 1.05, 3.15 + i * 0.55, 0.16, GREEN)
    text(s, 1.05, 3.12 + i * 0.55, 0.16, 0.16, [{"runs": [("✓", 10, True, WHITE, HEAD)], "align": PP_ALIGN.CENTER, "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.4, 3.06 + i * 0.55, 10.8, 0.5, [{"runs": [(t, 15, False, INK, BODY)], "sa": 0}])
rrect(s, 1.0, 5.05, 11.33, 0.95, RGBColor(0xEE,0xF3,0xFB), line=BLUE, lw=1.2, radius=0.1)
text(s, 1.3, 5.05, 10.8, 0.95, [{"runs": [("Next steps: ", 14, True, BLUE, HEAD), ("dbt · Apache Airflow · Cloud DW (BigQuery/Synapse) · Power BI Advanced · CI/CD", 13.5, False, INK, BODY)], "sa": 0}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.0, 6.3, 6.0, 0.5, [{"runs": [("ขอบคุณ — Q & A", 18, True, NAVY, HEAD)], "sa": 0}])
link_pill(s, 7.55, 6.34, 4.78, 0.52, "🔗 amornpan/advanced-de-applied-analytics", sz=11)

# page numbers — bottom-center, navy on the light area; skip the dark cover (slide 1)
_all = list(prs.slides)
for _i, _s in enumerate(_all, start=1):
    if _i == 1:
        continue
    text(_s, 6.0, 7.16, 1.33, 0.3,
         [{"runs": [("%d / %d" % (_i, len(_all)), 9.5, False, NAVY, BODY)], "align": PP_ALIGN.CENTER, "sa": 0}])

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
